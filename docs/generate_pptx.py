import os
import sys
import subprocess

# Ensure python-pptx is installed
try:
    import pptx
except ImportError:
    print("python-pptx is not installed. Installing it now...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    
    # Set to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme Colors
    NAVY = RGBColor(15, 44, 89)       # #0F2C59 - Primary Headers / Titles
    TEAL = RGBColor(0, 188, 212)      # #00BCD4 - Accents / Highlights
    DARK_TEXT = RGBColor(33, 33, 33)   # #212121 - Body Text
    WHITE = RGBColor(255, 255, 255)   # #FFFFFF - Backgrounds / Text on Navy
    LIGHT_BG = RGBColor(245, 245, 247) # #F5F5F7 - Slide Background
    MUTED_GRAY = RGBColor(117, 117, 117) # #757575 - Captions
    ACCENT_RED = RGBColor(211, 47, 47)  # #D32F2F - Danger / Problem highlight
    ACCENT_GREEN = RGBColor(56, 142, 60) # #388E3C - Success metrics
    
    # Blank layout for custom shapes positioning
    blank_layout = prs.slide_layouts[6]
    
    # --- HELPER FUNCTIONS ---
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
        
    def add_shape(slide, shape_type, left, top, width, height, fill_color, line_color=None):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape
        
    def add_text(slide, left, top, width, height, text, font_size=12, bold=False, color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name="Arial"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        return txBox

    def add_multiline_text(slide, left, top, width, height, paragraphs_data, font_name="Arial"):
        """
        paragraphs_data is a list of tuples: (text, font_size, bold, color, bullet, spacing_after)
        """
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for idx, data in enumerate(paragraphs_data):
            text, size, bold, color, bullet, space = data
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            if bullet:
                p.level = 0
            
            run = p.add_run()
            run.text = text
            run.font.name = font_name
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            if space > 0:
                p.space_after = Pt(space)
        return txBox

    def add_header(slide, title, category="PROJECT FLOWISE"):
        # Header category tag
        add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.3), category.upper(), font_size=10, bold=True, color=TEAL)
        # Header title
        add_text(slide, Inches(0.8), Inches(0.6), Inches(11.5), Inches(0.8), title, font_size=28, bold=True, color=NAVY)
        # Decorative divider line
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.3), Inches(11.733), Inches(0.04), TEAL)

    # =========================================================================
    # SLIDE 1: Title Slide (Dark Theme)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, NAVY)
    
    # Decorative accent block on left
    add_shape(slide1, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5), TEAL)
    
    # Project Title
    add_text(slide1, Inches(1.0), Inches(1.8), Inches(11), Inches(1.2), "FLOWISE", font_size=64, bold=True, color=WHITE)
    # Subtitle
    add_text(slide1, Inches(1.0), Inches(2.9), Inches(11), Inches(0.6), "Real-Time IoT & Machine Learning Water Leak Detection System", font_size=20, bold=False, color=TEAL)
    
    # Submission Metadata
    meta_p = [
        ("Technical Funding Proposal for Physical Pipeline Test-Bed", 14, True, WHITE, False, 20),
        ("Project Team:", 12, True, TEAL, False, 5),
        ("Sara Khan (NUML-F22-31354)  |  Batool Tariq (NUML-F22-16916)", 13, False, WHITE, False, 15),
        ("Supervised by:", 12, True, TEAL, False, 5),
        ("Ms. Qurat-ul-Ain Raja", 13, False, WHITE, False, 20),
        ("Department of Computer Science  |  Faculty of Engineering & CS", 11, False, MUTED_GRAY, False, 2),
        ("National University of Modern Languages (NUML), Islamabad", 11, False, MUTED_GRAY, False, 0)
    ]
    add_multiline_text(slide1, Inches(1.0), Inches(3.8), Inches(11.5), Inches(3.0), meta_p)

    # =========================================================================
    # SLIDE 2: Problem Statement & Significance
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, LIGHT_BG)
    add_header(slide2, "The Water Scarcity & Infrastructure Challenge")
    
    # Left Card (Key Metric)
    add_shape(slide2, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(4.5), Inches(5.1), WHITE)
    add_text(slide2, Inches(1.1), Inches(2.2), Inches(3.9), Inches(0.4), "NATIONAL EMERGENCY", font_size=11, bold=True, color=MUTED_GRAY)
    add_text(slide2, Inches(1.1), Inches(2.5), Inches(3.9), Inches(1.2), "30% - 40%", font_size=64, bold=True, color=ACCENT_RED)
    add_text(slide2, Inches(1.1), Inches(3.7), Inches(3.9), Inches(0.8), "Of municipal water is lost as Non-Revenue Water (NRW)", font_size=16, bold=True, color=NAVY)
    
    challenge_meta = [
        ("Pakistan is approaching absolute water scarcity.", 13, False, DARK_TEXT, False, 8),
        ("Pipes waste millions of gallons daily due to delayed, manual, and reactive leak reporting.", 13, False, DARK_TEXT, False, 0)
    ]
    add_multiline_text(slide2, Inches(1.1), Inches(4.7), Inches(3.9), Inches(1.8), challenge_meta)
    
    # Right Card (Traditional vs. Flowise)
    add_shape(slide2, MSO_SHAPE.RECTANGLE, Inches(5.6), Inches(1.6), Inches(6.933), Inches(5.1), WHITE)
    
    right_meta = [
        ("Traditional Detection Drawbacks", 18, True, NAVY, False, 12),
        ("• Manual Inspection: Slow, labor-intensive, and highly error-prone.", 13, False, DARK_TEXT, False, 6),
        ("• Delayed Reporting: Leaks detected only after visible surface flooding.", 13, False, DARK_TEXT, False, 6),
        ("• Water Contamination: Low-pressure pipe cracks draw in sewage contaminants.", 13, False, DARK_TEXT, False, 18),
        
        ("The Flowise Automation Advantage", 18, True, TEAL, False, 12),
        ("• Automated & Continuous: 24/7 edge-to-cloud monitoring.", 13, False, DARK_TEXT, False, 6),
        ("• High-Latency Warnings: Triggers alerts to managers in under 10 seconds.", 13, False, DARK_TEXT, False, 6),
        ("• Preventative Insights: Detects minor pressure drops before major pipe bursts.", 13, False, DARK_TEXT, False, 0)
    ]
    add_multiline_text(slide2, Inches(6.0), Inches(1.9), Inches(6.1), Inches(4.6), right_meta)

    # =========================================================================
    # SLIDE 3: Project Objectives
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, LIGHT_BG)
    add_header(slide3, "Core Technical Objectives")
    
    # Grid of 5 Cards
    card_width = Inches(3.6)
    card_height = Inches(2.2)
    gap_x = Inches(0.4)
    gap_y = Inches(0.4)
    
    coords = [
        (Inches(0.8), Inches(1.8)),
        (Inches(0.8 + 3.6 + 0.4), Inches(1.8)),
        (Inches(0.8 + 2 * (3.6 + 0.4)), Inches(1.8)),
        (Inches(0.8 + 0.8), Inches(1.8 + 2.2 + 0.4)),
        (Inches(0.8 + 0.8 + 3.6 + 0.4), Inches(1.8 + 2.2 + 0.4))
    ]
    
    objectives = [
        ("TO DESIGN", "IoT Edge Telemetry Node", "Assemble high-frequency sensor acquisition modules on the controller to capture flow rate and pipeline pressure data at regular intervals."),
        ("TO DEVELOP", "Real-Time Data Pipeline", "Construct an data transmission model linking the sensing node to a centralized cloud database."),
        ("TO TRAIN", "Machine Learning Model", "Train a high-recall Machine Learning model using custom rolling-window feature engineering to isolate leak events from transient noise."),
        ("TO DEPLOY", "Mobile App & Cloud", "Deploy the model on a cloud server and create a cross-platform mobile application with live charts and instant notifications."),
        ("TO VALIDATE", "Physical Test-Bed Validation", "Construct a pressurized PVC pipeline prototype to verify system latency, noise resilience, and prediction accuracy in physical trials.")
    ]
    
    for i, obj in enumerate(objectives):
        cx, cy = coords[i]
        tag, title, desc = obj
        # Card background
        add_shape(slide3, MSO_SHAPE.RECTANGLE, cx, cy, card_width, card_height, WHITE)
        # Accent top line
        add_shape(slide3, MSO_SHAPE.RECTANGLE, cx, cy, card_width, Inches(0.12), TEAL if i % 2 == 0 else NAVY)
        
        card_content = [
            (tag, 11, True, TEAL if i % 2 == 0 else NAVY, False, 2),
            (title, 14, True, NAVY if i % 2 == 0 else DARK_TEXT, False, 6),
            (desc, 11, False, DARK_TEXT, False, 0)
        ]
        add_multiline_text(slide3, cx + Inches(0.2), cy + Inches(0.2), card_width - Inches(0.4), card_height - Inches(0.4), card_content)

    # =========================================================================
    # SLIDE 4: System Architecture
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, LIGHT_BG)
    add_header(slide4, "Flowise End-to-End Architecture")
    
    # 4 Steps horizontal flow
    step_width = Inches(2.6)
    step_height = Inches(4.5)
    step_gap = Inches(0.4)
    start_x = Inches(0.8)
    y_pos = Inches(1.8)
    
    steps = [
        ("01", "Edge Telemetry", "Sensing Node", "Reads flow rates and pipeline pressure from sensor modules.", "Edge Controller", "Periodic Sampling"),
        ("02", "Cloud Bridge", "Central Database", "Stores structured, timestamped telemetry data. Serves as active sync hub for cloud & mobile app.", "Cloud Database", "Telemetry Streams"),
        ("03", "ML Inference", "Cloud Server", "Cloud background worker runs Machine Learning model. Computes rolling trends and applies temporal voting.", "ML Classifier", "Real-Time Processing"),
        ("04", "User Alerting", "Mobile Dashboard", "Listens to data updates, displays interactive live charts, and triggers system alerts.", "Mobile App", "<10s Total Latency")
    ]
    
    for i, step in enumerate(steps):
        num, title, sub, desc, tech, rate = step
        x = start_x + i * (step_width + step_gap)
        
        # Step Background
        add_shape(slide4, MSO_SHAPE.RECTANGLE, x, y_pos, step_width, step_height, WHITE)
        # Top number strip
        add_shape(slide4, MSO_SHAPE.RECTANGLE, x, y_pos, step_width, Inches(0.5), NAVY)
        add_text(slide4, x + Inches(0.2), y_pos + Inches(0.1), step_width - Inches(0.4), Inches(0.3), num, font_size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        
        step_content = [
            (title, 15, True, NAVY, False, 2),
            (sub, 12, True, TEAL, False, 10),
            (desc, 11, False, DARK_TEXT, False, 18),
            ("Unit: " + tech, 10, True, MUTED_GRAY, False, 4),
            ("Process: " + rate, 10, False, MUTED_GRAY, False, 0)
        ]
        add_multiline_text(slide4, x + Inches(0.2), y_pos + Inches(0.7), step_width - Inches(0.4), step_height - Inches(0.8), step_content)
        
        # Connecting arrow (skip for last step)
        if i < 3:
            arrow_x = x + step_width + Inches(0.05)
            arrow_y = y_pos + Inches(2.0)
            add_text(slide4, arrow_x, arrow_y, step_gap - Inches(0.1), Inches(0.5), "➔", font_size=24, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 5: System Integration & Software (Combined Hardware, ML, and App)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, LIGHT_BG)
    add_header(slide5, "Full-System Integration & Software")
    
    # Left Details Card
    add_shape(slide5, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.1), WHITE)
    
    hw_details = [
        ("System Components Overview", 18, True, NAVY, False, 8),
        ("We have built and validated a fully integrated system across physical hardware, cloud analysis, and mobile monitoring.", 12, False, DARK_TEXT, False, 15),
        
        ("Key Features & Functionality:", 14, True, TEAL, False, 6),
        ("• Physical Testbed (Hardware): Closed-loop pressurized PVC pipeline rig with flow sensors, pressure transmitters, and rechargeable battery power management to capture real-world data.", 11.5, False, DARK_TEXT, False, 10),
        ("• Leak Classifier (Machine Learning): Real-time analysis of flow differences and pressure changes using an intelligent model, stabilized by a voting filter to prevent false alarms.", 11.5, False, DARK_TEXT, False, 10),
        ("• User Dashboard (Mobile Application): Live telemetry visualization, secure user login, daily consumption logging, and instant push notification alerts when leakages occur.", 11.5, False, DARK_TEXT, False, 0)
    ]
    add_multiline_text(slide5, Inches(1.1), Inches(1.8), Inches(5.4), Inches(4.7), hw_details)
    
    # Right Image Card
    add_shape(slide5, MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(1.6), Inches(5.333), Inches(5.1), WHITE)
    
    img_path = "/home/sara/Desktop/Flowise/hardwaredev/FullHardwarePrototype.jpg"
    if os.path.exists(img_path):
        try:
            # Place image slightly inside the card boundaries
            slide5.shapes.add_picture(img_path, Inches(7.4), Inches(1.8), Inches(4.933), Inches(4.3))
            add_text(slide5, Inches(7.4), Inches(6.25), Inches(4.933), Inches(0.3), "Fig: Labeled Flowise Hardware Prototype", font_size=10, bold=False, color=MUTED_GRAY, align=PP_ALIGN.CENTER)
        except Exception as e:
            print(f"Error placing hardware image: {e}")
            add_text(slide5, Inches(7.4), Inches(3.5), Inches(4.933), Inches(1.0), "[ Image: Labeled Hardware Prototype ]\n(Failed to render image file)", font_size=14, bold=True, color=MUTED_GRAY, align=PP_ALIGN.CENTER)
    else:
        add_text(slide5, Inches(7.4), Inches(3.5), Inches(4.933), Inches(1.0), "[ Image: Labeled Hardware Prototype ]\n(FullHardwarePrototype.jpg not found)", font_size=14, bold=True, color=MUTED_GRAY, align=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 6: Cost Breakdown & Funding Request (Previous Slide 8)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, LIGHT_BG)
    add_header(slide6, "Technical Funding Request & Cost Breakdown")
    
    # Left Card: The Justification
    add_shape(slide6, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.0), Inches(5.1), WHITE)
    add_shape(slide6, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.0), Inches(0.12), NAVY)
    
    just_meta = [
        ("Funding Justification", 18, True, NAVY, False, 12),
        ("• Financial Constraints: As a two-student team, self-funding the total prototype development has cost 30,000 PKR, representing a heavy burden of 15,000 PKR per student.", 12.5, False, DARK_TEXT, False, 10),
        ("• Grant Application: Requesting a technical grant of 25,050 PKR from the university to offset hardware expenses.", 12.5, False, DARK_TEXT, False, 10),
        ("• Academic Return: Provides the Computer Science department with a robust, physical IoT-ML water monitoring apparatus for future research projects.", 12.5, False, DARK_TEXT, False, 15),
        
        ("Total Grant Request:", 14, True, TEAL, False, 4),
        ("25,050 PKR", 32, True, NAVY, False, 0)
    ]
    add_multiline_text(slide6, Inches(1.1), Inches(1.9), Inches(4.4), Inches(4.5), just_meta)
    
    # Right: Table Shape
    table_shape = slide6.shapes.add_table(8, 2, Inches(6.2), Inches(1.6), Inches(6.333), Inches(5.1))
    table = table_shape.table
    
    # Set Column Widths
    table.columns[0].width = Inches(4.5)
    table.columns[1].width = Inches(1.833)
    
    budget_items = [
        ("Component Category / Items", "Cost (PKR)"),
        ("Pipeline Infrastructure & Acrylic Mounting Base", "9,980"),
        ("Sensing Module (Flow & Pressure Sensors)", "7,900"),
        ("IoT & Processing Unit (Microcontroller, Battery, BMS)", "2,120"),
        ("Signal Conditioning & Voltage Regulators", "1,250"),
        ("Circuitry & Auxiliary Hardware (Wires, Switches)", "1,516"),
        ("Fluid Dynamics Hardware (12V Water Pump, Adapter)", "760"),
        ("Logistics (Shipping and Cumulative Taxes)", "1,524")
    ]
    
    for row_idx, item in enumerate(budget_items):
        desc, cost = item
        for col_idx, text in enumerate([desc, cost]):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            # Format text in cell
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
            p.font.name = "Arial"
            p.font.size = Pt(10.5 if row_idx > 0 else 12)
            p.font.bold = (row_idx == 0)
            p.font.color.rgb = WHITE if row_idx == 0 else DARK_TEXT
            
            # Formatting Row backgrounds
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 1 else RGBColor(235, 237, 240)

    # =========================================================================
    # SLIDE 7: Expected Deliverables & Social Impact (Previous Slide 9)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, LIGHT_BG)
    add_header(slide7, "Expected Deliverables & Project Impact")
    
    # Left Card: Deliverables
    add_shape(slide7, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1), WHITE)
    add_shape(slide7, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(0.12), NAVY)
    
    deliv_meta = [
        ("Project Deliverables", 18, True, NAVY, False, 12),
        ("1. Calibrated Physical Prototype: A functioning pressurized pipe testbed with flow and pressure logging.", 12.5, False, DARK_TEXT, False, 10),
        ("2. Cloud Inference Pipeline: Live Machine Learning inference running 24/7 on a cloud server.", 12.5, False, DARK_TEXT, False, 10),
        ("3. Mobile Application: User-ready mobile app with live charting and real-time alert notifications.", 12.5, False, DARK_TEXT, False, 10),
        ("4. Labeled Physical Dataset: Clean flow and pressure training dataset representing standard and simulated leak states.", 12.5, False, DARK_TEXT, False, 0)
    ]
    add_multiline_text(slide7, Inches(1.1), Inches(1.9), Inches(5.0), Inches(4.5), deliv_meta)
    
    # Right Card: Social Impact
    add_shape(slide7, MSO_SHAPE.RECTANGLE, Inches(6.933), Inches(1.6), Inches(5.6), Inches(5.1), WHITE)
    add_shape(slide7, MSO_SHAPE.RECTANGLE, Inches(6.933), Inches(1.6), Inches(5.6), Inches(0.12), TEAL)
    
    impact_meta = [
        ("Socio-Economic Impact", 18, True, NAVY, False, 12),
        ("• Combating Water Scarcity: Directly targets the 30-40% Non-Revenue Water loss in Pakistani municipal pipes.", 12.5, False, DARK_TEXT, False, 10),
        ("• Infrastructure Safety: Immediate alert under 10s prevents soil erosion, road damage, and building foundation risks.", 12.5, False, DARK_TEXT, False, 10),
        ("• Health Hazard Prevention: Prevents drinking water contamination caused by sewage seepage during pipe pressure drops.", 12.5, False, DARK_TEXT, False, 10),
        ("• Public Utility Optimization: Transition from slow manual reporting to automated digital audits, saving millions in repair overheads.", 12.5, False, DARK_TEXT, False, 0)
    ]
    add_multiline_text(slide7, Inches(7.233), Inches(1.9), Inches(5.0), Inches(4.5), impact_meta)

    # =========================================================================
    # SLIDE 8: Future Scalability & Commercialization (Previous Slide 10)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, LIGHT_BG)
    add_header(slide8, "Future Scalability & Commercialization")
    
    # Grid of 4 Cards
    col_w = Inches(5.6)
    row_h = Inches(2.35)
    
    # Coordinates
    x1, y1 = Inches(0.8), Inches(1.6)
    x2, y1 = Inches(6.933), Inches(1.6)
    x1, y2 = Inches(0.8), Inches(4.35)
    x2, y2 = Inches(6.933), Inches(4.35)
    
    scalability_items = [
        (x1, y1, "1. Multi-Node Leak Localization", "Upgrade the single-segment prototype to a multi-node mesh network. By analyzing pressure differential curves and flow differences across nodes, the system can localize the exact leak coordinates within a larger network."),
        (x2, y1, "2. Municipal Enterprise Dashboards", "Develop centralized platforms for municipal bodies (like WASA, KWSB, or CDA) to monitor city mains, automate leak alerts, and audit water distribution metrics."),
        (x1, y2, "3. Smart Housing Society Integration", "Market the solution as a smart infrastructure amenity for private housing developers (e.g., DHA, Bahria Town) to audit internal water billing and minimize leak damages."),
        (x2, y2, "4. Edge Intelligence", "Quantize the Machine Learning classifier to run directly on the edge controller. This enables localized inference, reducing connectivity dependence, latency, and cloud database costs.")
    ]
    
    for x, y, title, desc in scalability_items:
        add_shape(slide8, MSO_SHAPE.RECTANGLE, x, y, col_w, row_h, WHITE)
        add_shape(slide8, MSO_SHAPE.RECTANGLE, x, y, col_w, Inches(0.08), TEAL)
        
        card_content = [
            (title, 14, True, NAVY, False, 4),
            (desc, 11, False, DARK_TEXT, False, 0)
        ]
        add_multiline_text(slide8, x + Inches(0.2), y + Inches(0.15), col_w - Inches(0.4), row_h - Inches(0.2), card_content)

    # =========================================================================
    # SLIDE 9: Conclusion & Final Appeal (Previous Slide 11) (Dark Theme)
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9, NAVY)
    add_shape(slide9, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5), TEAL)
    
    add_text(slide9, Inches(1.0), Inches(0.8), Inches(11.3), Inches(0.4), "CONCLUSION & FINAL APPEAL", font_size=11, bold=True, color=TEAL)
    add_text(slide9, Inches(1.0), Inches(1.2), Inches(11.3), Inches(0.8), "Flowise: Conserving Pakistan's Water Reserves", font_size=28, bold=True, color=WHITE)
    add_shape(slide9, MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.1), Inches(11.3), Inches(0.04), TEAL)
    
    concl_meta = [
        ("• End-to-End System Readiness: The software core (machine learning models, database, and mobile application) is fully functional and ready.", 14, False, WHITE, False, 10),
        ("• Physical Prototype Necessity: Deploying on a physical pipeline is critical to calibrate our models against real-world environmental noise and sensor variance.", 14, False, WHITE, False, 10),
        ("• Financial Feasibility: Departmental funding of 25,050 PKR will cover physical sensor and pipeline overheads, making this student project viable.", 14, False, WHITE, False, 10),
        ("• Social Value: Demonstrates a low-cost, smart water leak detection model optimized for public utility infrastructure in Pakistan.", 14, False, WHITE, False, 0)
    ]
    add_multiline_text(slide9, Inches(1.0), Inches(2.4), Inches(11.3), Inches(2.2), concl_meta)
    
    # Endorsement and signatures block
    sig_p = [
        ("Sara Khan  |  Batool Tariq", 12, True, TEAL, False, 2),
        ("FYP Project Team, BSCS 8-A (Morning)", 11, False, WHITE, False, 18),
        ("Endorsed by Supervisor:", 12, True, TEAL, False, 2),
        ("Ms. Qurat-ul-Ain Raja", 12, True, WHITE, False, 2),
        ("Department of Computer Science, NUML, Islamabad", 11, False, MUTED_GRAY, False, 0)
    ]
    add_multiline_text(slide9, Inches(1.0), Inches(5.1), Inches(11.3), Inches(2.0), sig_p)
    
    # Save the presentation
    output_filename = "/home/sara/Desktop/Flowise/docs/Flowise_Funding_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation generated successfully: {output_filename}")

if __name__ == "__main__":
    build_presentation()
